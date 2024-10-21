#Código 100% funcional

import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

altura_slide = 7.5 # Altura do slide em polegadas
pos_x = 2.65     # Posição X
pos_y = 2     # Posição Y
largura_textbox = 8   # Largura textbox
altura_textbox = 2    # Altura textbox
font_size = 48  # Tamanho da fonte
font_name = "BANDEX" # Fonte
cor = [255, 255, 255]


def selecionar_arquivo():
    """Abre um seletor de arquivo para escolher arquivos .txt."""
    root = tk.Tk()
    root.withdraw()  # Esconde a janela principal do tkinter
    arquivos_txt = filedialog.askopenfilenames(
        title="Selecione os arquivos .txt",
        filetypes=[("Arquivo de Texto", "*.txt")]
    )
    return arquivos_txt  # Retorna uma tupla com os caminhos dos arquivos

def calcular_polegadas(EMU):
    return EMU / 914400

def calcular_largura(altura):
    """Calcula a largura com base na altura para manter a proporção 16:9."""
    return (16 / 9) * altura
    

def criar_apresentacao(txt_file, pos_x, pos_y, largura_textbox, altura_textbox, font_size, font_name, r, g, b):
    """Cria uma apresentação PowerPoint a partir de um arquivo .txt."""
    # Define o nome do arquivo PowerPoint de saída com base no nome do arquivo de texto
    arquivo_ppt = f'{txt_file.split("/")[-1].replace(".txt", "")}.pptx'
    
    prs = Presentation()

    largura_slide = calcular_largura(altura_slide)
    # Define o tamanho da apresentação
    prs.slide_width = Inches(largura_slide)
    prs.slide_height = Inches(altura_slide)

    
    print("A largura e altura da apresentação é: ", prs.slide_width / 914400, prs.slide_height)

    # Abre o arquivo .txt e lê as linhas
    with open(txt_file, 'r', encoding='utf-8') as file:
        linhas = file.readlines()

    path_img = add_imagem()
    print("path_img: ", path_img)


    titulo = True
    # Para cada linha do arquivo .txt, cria um slide e adiciona o texto
    for indice, linha in enumerate(linhas):
        linha = linha.strip()  # Remove espaços e quebras de linha adicionais
        # desativado pra perfomance: print(linha)

        if linha:  # Só adiciona se a linha não for vazia
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Cria um slide branco e sem título

            largura_polegadas = prs.slide_width
            altura_polegadas = prs.slide_height

            # Verifica se o slide tem um título
            if slide.shapes.title:
                # Verifica se o título está vazio
                if not slide.shapes.title.text.strip():
                    # Remove a caixa de título se estiver vazia
                    slide.shapes._spTree.remove(slide.shapes.title._element)

            # Adiciona a imagem ajustada ao tamanho do slide
            slide.shapes.add_picture(path_img, 0, 0, largura_polegadas, altura_polegadas)

            textbox = slide.shapes.add_textbox(Inches(pos_x), Inches(pos_y), Inches(largura_textbox), Inches(altura_textbox)) # Adiciona o textbox
            text_frame = textbox.text_frame                 # Acessa o textbox
            text_frame.text = linha                         # Insere o conteúdo da string linha no textbox
            text_frame.auto_size = True                     # Ajusta automaticamente o tamanho da caixa
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # Centraliza verticalmente

            p = text_frame.paragraphs[0]         # Acessa o primeiro (e único) parágrafo
            p.font.size = Pt(font_size)          # Define o tamanho da fonte
            p.font.name = font_name              # Define a fonte
            p.font.color.rgb = RGBColor(r, g, b) # Define a cor
            p.alignment = PP_ALIGN.CENTER


            if titulo:
                print("Adicionando titulo")
                p.font.size = Pt(font_size * 1.6)  # Define o tamanho da fonte
                titulo = False
            elif not titulo and len(linha)<40:
                print(f"Linha: {indice}")
            elif not titulo and len(linha)>39:
                print(f"Linha grande demais:  {indice} ")

    # Salva a apresentação PowerPoint
    try:
        prs.save(arquivo_ppt)
        print(f'Apresentação {arquivo_ppt} criada com sucesso!')
    except:
        print("Erro ao salvar o arquivo, verifique se há uma janela aberta no PowerPoint")
    input("Aperte enter pra sair")

def add_imagem():
    janela_add_img = tk.Toplevel()
    janela_add_img.withdraw()  # Esconde a janela principal do tkinter
    arquivo_img = filedialog.askopenfilenames(
        title="Selecione o arquivo .jpg ou png para fundo do slide",
        filetypes=[
            ("Arquivo de Imagem", "*.jpg"),
            ("Arquivo de Imagem", "*.png")]
    )
    
    if arquivo_img:
        return arquivo_img[0]  # Retorna apenas o primeiro arquivo selecionado
    else:
        return None  # Caso nenhum arquivo seja selecionado


r, g, b = cor


# Seleciona os arquivos .txt
arquivos_txt = selecionar_arquivo()

# Cria a apresentação PowerPoint
if arquivos_txt:
    
    for arquivo_txt in arquivos_txt:  # Itera sobre cada arquivo selecionado
        criar_apresentacao(arquivo_txt, pos_x, pos_y, largura_textbox, altura_textbox, font_size, font_name, r, g, b)
