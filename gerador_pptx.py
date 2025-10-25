#v 2.27
import tkinter as tk
from tkinter import messagebox, colorchooser, filedialog, ttk
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

version = "2.27"

# Inicializa a variável path_img como None para evitar erros de referência antes da seleção
path_img = None

def limpar():
    """Limpa o campo de texto."""
    global campo_texto, lista_arquivos
    if 'lista_arquivos' in globals() and lista_arquivos.winfo_exists():
        lista_arquivos.delete(0, tk.END)
    if 'campo_texto' in globals() and campo_texto.winfo_exists():
        campo_texto.delete(1.0, tk.END)  # Limpa todo o conteúdo do campo de texto

# Funções para a interface
def colar_conteudo():
    try:
        # Cola o conteúdo da área de transferência no campo de texto
        texto = janela.clipboard_get()
        campo_texto.delete(1.0, tk.END)  # Limpa o campo de texto
        campo_texto.insert(tk.END, texto)  # Insere o conteúdo colado
    except tk.TclError:
        messagebox.showerror("Erro", "Nada para colar!")

def selecionar_txt():
    # Abre o seletor de arquivos para escolher um arquivo de texto
    arquivos_txt = filedialog.askopenfilename(
        title="Selecionar Arquivo de Texto",
        filetypes=(("Arquivos de Texto", "*.txt"), ("Todos os arquivos", "*.*")))
    
    texto = ""

    if arquivos_txt:
        # for arquivo in arquivos_txt:
        with open(arquivos_txt, 'r', encoding='utf-8') as file:
            texto += file.read() + "\n"

            campo_texto.delete(1.0, tk.END)  # Limpa o campo de texto
            campo_texto.insert(tk.END, texto)  # Insere o conteúdo colado        
    else:
        messagebox.showwarning("Aviso", "Nenhum arquivo selecionado")

def selecionar_arquivos():
    global campo_texto, lista_arquivos, arquivos

    arquivos = filedialog.askopenfilenames(
        title="Selecionar Arquivo(s) de Texto",
        filetypes=(("Arquivos de Texto", "*.txt"), ("Todos os arquivos", "*.*"))
    )

    if not arquivos:
        messagebox.showwarning("Aviso", "Nenhum arquivo selecionado")
        return

    # Se houver somente um arquivo, exibe o conteúdo no campo de texto
    if len(arquivos) == 1:
        # Se houver uma Listbox visível, destruir
        if 'lista_arquivos' in globals() and lista_arquivos.winfo_exists():
            lista_arquivos.destroy()

        # Exibir o Text
        campo_texto.grid(row=1, column=0, columnspan=10, pady=5)

        with open(arquivos[0], 'r', encoding='utf-8') as file:
            conteudo = file.read()

        campo_texto.delete(1.0, tk.END)
        campo_texto.insert(tk.END, conteudo)

    # Se houver mais de um arquivo, exibe os nomes na Listbox
    else:
        # Se houver o Textbox visível, destruir
        if campo_texto.winfo_exists():
            campo_texto.grid_forget()

        # Se já existe Listbox, limpa. Senão, cria.
        if 'lista_arquivos' in globals() and lista_arquivos.winfo_exists():
            lista_arquivos.delete(0, tk.END)
        else:
            lista_arquivos = tk.Listbox(frame_texto, width=40, height=20)
            lista_arquivos.grid(row=1, column=0, columnspan=10, pady=5)

        # Insere os nomes dos arquivos na Lista
        for arq in arquivos:
            lista_arquivos.insert(tk.END, arq.split('/')[-1])

def selecionar_imagem():
    # Abre o seletor de arquivos para escolher uma imagem
    global path_img
    path_img = filedialog.askopenfilename(
    title="Selecionar Imagem",
    filetypes=(("Arquivos de Imagem", "*.png;*.jpg;*.jpeg"), ("Todos os arquivos", "*.*")))
    if path_img:
        label_caminho_imagem.config(text=f"Imagem: {path_img}")
    else:
        label_caminho_imagem.config(text="Nenhuma imagem selecionada")

janela_ajuda_aberta = False

def fechar_ajuda():
    # Fecha a janela de ajuda se estiver aberta
    global janela_ajuda_aberta, janela
    if janela_ajuda_aberta:
        janela_ajuda.destroy()
        janela_ajuda_aberta = False

def exibir_ajuda():
    global janela_ajuda_aberta, janela_ajuda
    # Verifica se a janela de ajuda já está aberta
    if janela_ajuda_aberta:
        return
    
    
    # Definições da janela de ajuda
    janela_ajuda = tk.Toplevel(janela)
    janela_ajuda.title("Ajuda")
    janela_ajuda.geometry("400x300")
    texto_ajuda = """
    Como usar o Gerador de Apresentações PPTX:

    1. Cole o conteúdo dos slides no campo de texto.
    2. Ajuste as configurações de layout e estilo.
    3. Selecione uma imagem de fundo (opcional).
    4. Clique em "Gerar PPTX" para criar a apresentação.

    Dicas:
    - Use o botão "Colar" para colar texto da área de transferência.
    - Escolha uma cor para o texto usando o seletor de cores.
    - Insira o nome da fonte e o tamanho desejados.
    """
    janela_ajuda.protocol("WM_DELETE_WINDOW", ao_fechar_ajuda)  # Fecha a janela de ajuda corretamente

    # Cria uma nova janela com instruções de uso
    label_ajuda = tk.Label(janela_ajuda, text=texto_ajuda, justify=tk.LEFT)
    label_ajuda.pack(padx=10, pady=10)
    janela_ajuda_aberta = True

def ao_fechar_ajuda():
    # Fecha a janela de ajuda e atualiza a variável de controle
    global janela_ajuda_aberta, janela_ajuda
    janela_ajuda_aberta = False
    janela_ajuda.destroy()

def calcular_largura(altura):
    """Calcula a largura com base na altura para manter a proporção 16:9."""
    return (16 / 9) * altura

def hex_to_rgb(hex_color):
    """Converte uma cor hexadecimal para RGB."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

# Seletor de cor
def escolher_cor():
    cor = colorchooser.askcolor()[1]  # Retorna o código hexadecimal da cor
    if cor:
        cor_selecionada.set(cor)  # Armazena a cor selecionada
        frame_cor.config(bg=cor)  # Atualiza a cor do frame
        # botao_cor.config(bg=cor)

def processar_arquivo_unico():
    """Cria uma apresentação PowerPoint a partir do conteúdo do campo de texto."""
    global path_img
    # Verifica se há imagem de fundo selecionada
    if not path_img:
        processar_sem_imagem = messagebox.askyesno("Aviso", "Nenhuma imagem de fundo selecionada. Deseja continuar sem imagem?")
        if not processar_sem_imagem:
            return

    if 'lista_arquivos' in globals() and lista_arquivos.winfo_exists():
        processar_arquivo_multiplo()

    # Coleta as configurações da interface
    altura_slide = float(7.5)  # Altura padrão
    font_size = int(font_size_entry.get())
    font_name = font_name_entry.get()
    n_maximo = int(n_maximo_entry.get())
    r, g, b = hex_to_rgb(cor_selecionada.get())  # Converte a cor hexadecimal para RGB

    # Coleta o conteúdo do campo de texto
    linhas = campo_texto.get(1.0, tk.END).splitlines()
    if is_maiusculas.get():
        linhas.upper()  # Converte todo o texto para maiúsculas
    if not linhas:
        messagebox.showwarning("Aviso", "O campo de texto está vazio!")
        return

    # Define o nome do arquivo com base na primeira linha do texto
    nome_arquivo = linhas[0].strip() + ".pptx"
    arquivo_ppt = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("Arquivos PowerPoint", "*.pptx")], initialfile=nome_arquivo)
    if not arquivo_ppt:
        return  # Usuário cancelou a operação

    # Cria a apresentação
    prs = Presentation()
    largura_slide = calcular_largura(altura_slide)
    prs.slide_width = Inches(largura_slide)
    prs.slide_height = Inches(altura_slide)

    # Calcula a posição do textbox considerando a fração de 10
    pos_x = prs.slide_width / 10 * int(pos_x_entry.get()) # Calcula a posição horizontal baseada na entrada do usuário
    pos_y = prs.slide_height / 10 * int(pos_y_entry.get()) # Calcula a posição vertical baseada na entrada do usuário

    # Verifica linhas grandes
    linhas_grandes = [i + 1 for i, linha in enumerate(linhas) if len(linha) > n_maximo]

    titulo = True
    for indice, linha in enumerate(linhas):
        linha = linha.rstrip()  # Remove espaços adicionais
        if not linha:
            continue  # Ignora linhas vazias

        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Cria um slide branco e sem título

        # Adiciona a imagem de fundo, se selecionada
        if path_img:
            slide.shapes.add_picture(path_img, 0, 0, prs.slide_width, prs.slide_height)
        else:
            pass  # Se não houver imagem, não faz nada

        # Adiciona o textbox
        textbox = slide.shapes.add_textbox(pos_x, pos_y, 8, 2) # Valores 8 e 2 não importam, pois o tamanho é ajustado automaticamente
        text_frame = textbox.text_frame
        text_frame.text = linha
        text_frame.auto_size = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = text_frame.paragraphs[0]
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.font.color.rgb = RGBColor(r, g, b)
        p.alignment = PP_ALIGN.CENTER

        if titulo:
            p.font.size = Pt(font_size * 1.6)  # Aumenta o tamanho da fonte do título
            titulo = False

        # Adiciona um comentário no primeiro slide com as linhas grandes
        if indice == 0 and linhas_grandes:
            comentario = f"Linhas grandes (>{n_maximo} caracteres): {', '.join(map(str, linhas_grandes))}"
            slide.notes_slide.notes_text_frame.text = comentario

    # Salva a apresentação
    prs.save(arquivo_ppt)
    messagebox.showinfo("Sucesso", f"Apresentação salva como {arquivo_ppt}")
            # ----- Fim da função criar_apresentacao -----

def processar_arquivo_multiplo():

    global arquivos
    """Cria uma apresentação PowerPoint para cada arquivo de texto."""
    global path_img

    try:
        arquivos_selecionados = lista_arquivos.get(0, tk.END)
    except AttributeError:
        print("Nenhum arquivo selecionado na lista.")

    if not arquivos_selecionados:
        messagebox.showwarning("Aviso", "Nenhum arquivo selecionado!")
        return
    # Verifica se há imagem de fundo selecionada
    if not path_img:
        processar_sem_imagem = messagebox.askyesno("Aviso", "Nenhuma imagem de fundo selecionada. Deseja continuar sem imagem?")
        if not processar_sem_imagem:
            return
        
    # Coleta as configurações da interface
    altura_slide = float(7.5)  # Altura padrão
    font_size = int(font_size_entry.get())
    font_name = font_name_entry.get()
    n_maximo = int(n_maximo_entry.get())
    r, g, b = hex_to_rgb(cor_selecionada.get())  # Converte a cor hexadecimal para RGB

    for index, nome_arquivo in enumerate(arquivos_selecionados):

        # Coleta o conteúdo do arquivo
        with open(arquivos[index], 'r', encoding='utf-8') as file:
            linhas = file.readlines()
        
        if is_maiusculas.get():
            linhas.upper()  # Converte todo o texto para maiúsculas

        # Define o nome do arquivo
        nome_arquivo = nome_arquivo.split('/')[-1].replace('.txt', '') + ".pptx"
        arquivo_ppt = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("Arquivos PowerPoint", "*.pptx")], initialfile=nome_arquivo)
        if not arquivo_ppt:
            return  # Usuário cancelou a operação

        # Cria a apresentação
        prs = Presentation()
        largura_slide = calcular_largura(altura_slide)
        prs.slide_width = Inches(largura_slide)
        prs.slide_height = Inches(altura_slide)

        # Calcula a posição do textbox considerando a fração de 10
        pos_x = prs.slide_width / 10 * int(pos_x_entry.get())  # Calcula a posição horizontal baseada na entrada do usuário
        pos_y = prs.slide_height / 10 * int(pos_y_entry.get())  # Calcula a posição vertical baseada na entrada do usuário

        # Verifica linhas grandes
        linhas_grandes = [i + 1 for i, linha in enumerate(linhas) if len(linha) > n_maximo]

        titulo = True

        for indice, linha in enumerate(linhas):
            linha = linha.rstrip()  # Remove espaços adicionais
            # if not linha:
            #     continue  # Ignora linhas vazias
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Cria um slide branco e sem título

            # Adiciona a imagem de fundo, se selecionada
            if path_img:
                slide.shapes.add_picture(path_img, 0, 0, prs.slide_width, prs.slide_height)
            else:
                pass  # Se não houver imagem, não faz nada

            # Adiciona o textbox
            textbox = slide.shapes.add_textbox(pos_x, pos_y, 8, 2)  # Valores 8 e 2 não importam, pois o tamanho é ajustado automaticamente
            textbox = slide.shapes.add_textbox(pos_x, pos_y, 8, 2) # Valores 8 e 2 não importam, pois o tamanho é ajustado automaticamente
            text_frame = textbox.text_frame
            text_frame.text = linha
            text_frame.auto_size = True
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            p = text_frame.paragraphs[0]
            p.font.size = Pt(font_size)
            p.font.name = font_name
            p.font.color.rgb = RGBColor(r, g, b)
            p.alignment = PP_ALIGN.CENTER

            if titulo:
                p.font.size = Pt(font_size * 1.6)  # Aumenta o tamanho da fonte do título
                titulo = False

            # Adiciona um comentário no primeiro slide com as linhas grandes
            if indice == 0 and linhas_grandes:
                comentario = f"Linhas grandes (>{n_maximo} caracteres): {', '.join(map(str, linhas_grandes))}"
                slide.notes_slide.notes_text_frame.text = comentario

        # Salva a apresentação
        prs.save(arquivo_ppt)
        messagebox.showinfo("Sucesso", f"Apresentação salva como {arquivo_ppt}")


# Configuração da janela principal
janela = tk.Tk()
janela.title("Gerador de Apresentações PPTX - " + version)
janela.geometry("620x500")

# --- Frame para o campo de texto ---
frame_texto = tk.Frame(janela)
frame_texto.grid(row=0, column=0, columnspan=2, padx=10, pady=10)

# Botão "Colar" acima do campo de texto
botao_colar = tk.Button(frame_texto, text="Colar", command=colar_conteudo)
botao_colar.grid(row=0, column=0,padx=5, pady=5, sticky="n")  # Alinha o botão à esquerda

# Botão selecionar arquivos de texto
botao_selecionar_txt = tk.Button(frame_texto, text="Selecionar txt", width=15, command=selecionar_arquivos)
botao_selecionar_txt.grid(row=0, column=1, padx=5, pady=5, sticky="n")  # Alinha o botão à esquerda

# Botão "Limpar" para limpar o campo de texto
botao_limpar = tk.Button(frame_texto, text="Limpar", command=limpar)
botao_limpar.grid(row=0, column=2, padx=10, pady=5, sticky="n")  # Alinha o botão à esquerda

# Campo de texto para o conteúdo dos slides
campo_texto = tk.Text(frame_texto, width=40, height=20)
campo_texto.grid(row=1, column=0, columnspan=10, pady=5)

# --- Frame para as definições ---
frame_definicoes = tk.Frame(janela)
frame_definicoes.grid(row=0, column=2, padx=10, pady=10, sticky="n")

# Notebook (TabControl) para organizar as configurações
notebook = ttk.Notebook(frame_definicoes)
notebook.grid(row=0, column=0, padx=10, pady=10, sticky="n")

# Aba "Formato"
aba_formato = ttk.Frame(notebook)
notebook.add(aba_formato, text="Formato")

# Controles de configuração na aba "Formato"

label_ajuda_formato = tk.Label(aba_formato,
text="Posicione o textbox usando valores de 0 a 10. \n Considere o centro da caixa de texto.")
label_ajuda_formato.grid(row=0, column=0, pady=5, padx=5, sticky="w") # Largura e altura do slide

# Label para definir a posição horizontal do textbox
tk.Label(aba_formato, text="Posição horizontal:").grid(row=2, column=0, sticky="w") # pos_x

# Entry para definir a posição horizontal do textbox
pos_x_entry = tk.Entry(aba_formato)
pos_x_entry.grid(row=3, column=0, pady=2, sticky="w")
pos_x_entry.insert(0, "5")  # Posição horizontal padrão
pos_x_entry.config(width=5)  # Define a largura do campo de entrada

# Label para definir a posição vertical do textbox
tk.Label(aba_formato, text="Posição vertical:").grid(row=4, column=0, sticky="w") # pos_y

# Entry para definir a posição vertical do textbox
pos_y_entry = tk.Entry(aba_formato)
pos_y_entry.grid(row=5, column=0, pady=2, sticky="w")
pos_y_entry.insert(0, "4")  # Posição vertical padrão
pos_y_entry.config(width=5)  # Define a largura do campo de entrada

# Labek para definir a quantidade de letras máxima antes de estourar o slide
tk.Label(aba_formato, text="Tamanho Máximo de Letras (nMaximo):").grid(row=10, column=0, sticky="w")
n_maximo_entry = tk.Entry(aba_formato, width=5)
n_maximo_entry.grid(row=11, column=0, pady=2, sticky="w")
n_maximo_entry.insert(0, "40")

is_maiusculas = tk.BooleanVar(value=True)  # Variável para controlar se o texto deve ser convertido para maiúsculas
# Checkbutton para converter texto para maiúsculas
check_maiusculas = tk.Checkbutton(aba_formato, text="Converter texto para MAIÚSCULAS", variable=is_maiusculas)
check_maiusculas.grid(row=12, column=0, pady=2, sticky="w")

# Aba "Texto"
aba_texto = ttk.Frame(notebook)
notebook.add(aba_texto, text="Texto")

# Controles de configuração na aba "Texto"
tk.Label(aba_texto, text="Tamanho da Fonte:").grid(row=0, column=0, sticky="w")
font_size_entry = tk.Entry(aba_texto)
font_size_entry.grid(row=1, column=0, pady=2)
font_size_entry.insert(0, "46")

# Label para o nome da fonte
tk.Label(aba_texto, text="Nome da Fonte:").grid(row=2, column=0, sticky="w")
font_name_entry = tk.Entry(aba_texto)
font_name_entry.grid(row=3, column=0, pady=2)
font_name_entry.insert(0, "BANDEX")

# Botão para selecionar a cor do texto
tk.Label(aba_texto, text="Cor do Texto:").grid(row=4, column=0, sticky="w")
botao_cor = tk.Button(aba_texto, text="Escolher Cor", command=escolher_cor, bg="white")
botao_cor.grid(row=5, column=0, pady=2)

cor_selecionada = tk.StringVar(value="#FFFFFF")  # inicialização da variável de cor com branco

# Cria um frame para exibir a cor selecionada
frame_cor = tk.Frame(aba_texto, width=20, height=20, bg=cor_selecionada.get())
frame_cor.grid(row=5, column=1, pady=2, sticky="w")

# Botão "Gerar PPTX"
botao_gerar = tk.Button(frame_definicoes, text="Gerar PPTX", width=20, height=3, command=processar_arquivo_unico)
botao_gerar.grid(row=1, column=0, padx=5, pady=15, sticky="n")

# Botão "Ajuda"
botao_ajuda = tk.Button(frame_definicoes, text="Ajuda", command=exibir_ajuda)
botao_ajuda.grid(row=2, column=0, padx=15, sticky="n")

# Define a aba "Texto" como padrão
notebook.select(aba_texto)

# --- Frame para os botões ---
frame_selecao_imagem = tk.Frame(janela)
frame_selecao_imagem.grid(row=2, column=0, columnspan=2, pady=10)

# Botão "Selecionar Imagem"
botao_imagem = tk.Button(frame_selecao_imagem, text="Selecionar Imagem", command=selecionar_imagem)
botao_imagem.grid(row=0, column=0, padx=5)

# Botão "Teste"
# botao_teste = tk.Button(frame_selecao_imagem, text="Teste", command=selecionar_arquivos)
# botao_teste.grid(row=0, column=1, padx=5)

# Rótulo para exibir o caminho da imagem
label_caminho_imagem = tk.Label(frame_selecao_imagem, text="Nenhuma imagem selecionada", fg="gray", wraplength=300)  # Define a largura máxima em pixels para o texto
label_caminho_imagem.grid(row=1, column=0, padx=5, pady=2)

# Inicia a interface
janela.mainloop()
